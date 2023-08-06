from pptx import Presentation


def find_and_replace_presentation(filename, old_word, new_word):
    """
    Replace instances of the given old word with the given new word in given PowerPoint file
    Based on Sam Redway's and Ricky Gonce's answers on :
        https://stackoverflow.com/questions/37924808/python-pptx-power-point-find-and-replace-text-ctrl-h
    :param filename: The name of the file to be edited
    :param old_word: The word to be searched for, assumed to be capitalized (e.g. 'Shark' not 'shark')
    :param new_word: The word to replace old_word with, assumed to be capitalized (e.g. 'Whale' not 'WHALE')
    :return: None
    """
    pres = Presentation(filename)
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        cur_text = run.text
                        new_text = find_and_replace(cur_text, old_word, new_word)
                        run.text = new_text
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        new_text = find_and_replace(cell.text, old_word, new_word)
                        cell.text = new_text
    pres.save(filename)


def find_and_replace(string, old_word, new_word):
    """
    Replace instances of the given old word with the given new word in the given string
    :param string: The string to be searched within
    :param old_word: The word to search for
    :param new_word: The word to replace instances of the old word with
    :return: The inputted string with instances of the old word replaced with instances of the new word
    """
    # old -> new
    lowercase_replaced = string.replace(old_word.lower(), new_word.lower())
    # OLD -> NEW
    uppercase_replaced = lowercase_replaced.replace(old_word.upper(), new_word.upper())
    # Old -> New
    capitalized_replaced = uppercase_replaced.replace(old_word, new_word)

    return capitalized_replaced


find_and_replace_presentation('TestFiles/shark.pptx', 'Whale', 'Shark')