from pptx import Presentation
from pptx import exc
import eventlog
from word import str_find_and_replace


def pptx_find_and_replace(filepath, old_word, new_word):
    """
    Replace instances of the given old word with the given new word in given PowerPoint file
    Based on Sam Redway's and Ricky Gonce's answers on :
        https://stackoverflow.com/questions/37924808/python-pptx-power-point-find-and-replace-text-ctrl-h
    :param filepath: The name of the file to be edited
    :param old_word: The word to be searched for, assumed to be capitalized (e.g. 'Shark' not 'shark')
    :param new_word: The word to replace old_word with, assumed to be capitalized (e.g. 'Whale' not 'WHALE')
    :return: None
    """
    try:
        pres = Presentation(filepath)

        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            cur_text = run.text
                            new_text = str_find_and_replace(cur_text, old_word, new_word)
                            run.text = new_text
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            new_text = str_find_and_replace(cell.text, old_word, new_word)
                            cell.text = new_text
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame
                notes.text = str_find_and_replace(notes.text, old_word, new_word)

        pres.save(filepath)
        eventlog.log_event("Replaced instances of " + old_word + " within" + filepath)
    except PermissionError:
        eventlog.log_event("ERROR: Could not edit contents of " + filepath + " because this file is already in use.")
    except exc.PackageNotFoundError:
        eventlog.log_event("ERROR: Could not edit contents of " + filepath + " because this file does not exist.")